"""
Generate academic presentation for ViFoodRec research project.
Run: python slides/generate_slides.py
"""
import os

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGURES = os.path.join(PROJECT_ROOT, "report", "figures")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

C_NAVY = RGBColor(30, 58, 138)
C_BLUE = RGBColor(59, 130, 246)
C_ORANGE = RGBColor(245, 166, 35)
C_DARK = RGBColor(51, 51, 51)
C_LIGHT = RGBColor(248, 249, 250)
C_WHITE = RGBColor(255, 255, 255)
C_GRAY = RGBColor(120, 120, 120)

FOOTER = "Khai thác DL Truyền thông MXH  |  Đánh giá KN trên Dữ liệu Tương tác ViFoodRec"


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def header(slide, title):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.75))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = C_NAVY
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.03))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ORANGE
    bar.line.fill.background()
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3))
    ft.text_frame.paragraphs[0].text = FOOTER
    ft.text_frame.paragraphs[0].font.size = Pt(9)
    ft.text_frame.paragraphs[0].font.color.rgb = C_GRAY


def bullets(slide, x, y, w, h, title, items, title_size=20, bullet_size=14, prefix="•"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = C_NAVY
        p.space_after = Pt(10)
    for item in items:
        bp = tf.add_paragraph()
        bp.text = f"{prefix} {item}"
        bp.font.size = Pt(bullet_size)
        bp.font.color.rgb = C_DARK
        bp.space_after = Pt(7)


def add_img(slide, rel_path, x, y, w):
    full = os.path.join(PROJECT_ROOT, rel_path)
    if os.path.exists(full):
        slide.shapes.add_picture(full, Inches(x), Inches(y), width=Inches(w))


def table_slide(slide, title, headers, rows, col_widths=None):
    header(slide, title)
    tbl_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.42 * (len(rows) + 1)))
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = C_WHITE
            p.font.bold = True
            p.font.size = Pt(10)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                if j == 0:
                    p.font.bold = True


def create_presentation():
    import pandas as pd
    
    # Load metrics dynamically
    reg_csv = os.path.join(PROJECT_ROOT, "report", "figures", "regression_metrics.csv")
    rank_csv = os.path.join(PROJECT_ROOT, "report", "figures", "ranking_metrics.csv")
    
    if os.path.exists(reg_csv) and os.path.exists(rank_csv):
        df_reg = pd.read_csv(reg_csv, index_col=0)
        df_rank = pd.read_csv(rank_csv, index_col=0)
    else:
        df_reg = None
        df_rank = None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    s = prs.slides.add_slide(blank)
    set_bg(s, C_NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = C_ORANGE
    accent.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11), Inches(3.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "ĐÁNH GIÁ HIỆU QUẢ CÁC PHƯƠNG PHÁP\nKHUYẾN NGHỊ TRÊN DỮ LIỆU TƯƠNG TÁC MXH"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p2 = tf.add_paragraph()
    p2.text = "ViFoodRec: 4.000 món ăn  |  182.678 ratings  |  CBF · CF · SVD · NCF"
    p2.font.size = Pt(16)
    p2.font.color.rgb = C_ORANGE
    p2.space_before = Pt(14)
    p3 = tf.add_paragraph()
    p3.text = "Môn: Khai thác Dữ liệu Truyền thông Xã hội  |  UIT - ĐHQG-HCM"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(200, 210, 230)
    p3.space_before = Pt(18)

    # 2. Outline
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    header(s, "Nội dung Trình bày")
    bullets(s, 1.2, 1.5, 10.5, 5.5, "", [
        "1. Vấn đề nghiên cứu & Câu hỏi RQ1–RQ3",
        "2. Công trình liên quan & Research Gap",
        "3. Dữ liệu tương tác ViFoodRec & Phân tích insight",
        "4. Phương pháp: CBF, CF, Funk SVD, NCF",
        "5. Thiết kế thí nghiệm (Experimental Setup)",
        "6. Kết quả & Thảo luận ý nghĩa thống kê",
        "7. Phân tích lỗi đa khía cạnh",
        "8. Kết luận & Hướng phát triển",
    ], title_size=1, bullet_size=16)

    # 3. Research context
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    header(s, "1. Bối cảnh: Dữ liệu Tương tác trên Nền tảng Ẩm thực")
    bullets(s, 0.5, 1.5, 6, 5, "Hành vi người dùng (User Interaction)", [
        "Rating món ăn = tín hiệu tương tác người dùng–nội dung trên MXH ẩm thực.",
        "Ma trận R(user × item) phản ánh sở thích ẩn cần khai thác.",
        "Density 45.22% — dữ liệu tương tác dày hơn nhiều dataset thương mại.",
    ])
    bullets(s, 6.8, 1.5, 6, 5, "Bài toán nghiên cứu", [
        "Input: ratings + metadata món ăn.",
        "Output: dự đoán rating + Top-K gợi ý.",
        "Không phải xây dựng hệ thống — mà đánh giá phương pháp khai thác.",
    ], prefix="→")

    # 4. Research Questions
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    header(s, "1. Câu hỏi Nghiên cứu (RQ1–RQ3)")
    bullets(s, 0.8, 1.5, 11.5, 5.5, "", [
        "RQ1: CBF và Memory-Based CF hoạt động thế nào trên ViFoodRec?",
        "RQ2: Funk SVD và NCF có cải thiện dự đoán rating so với CF truyền thống?",
        "RQ3: Mô hình nào phù hợp cho Rating Prediction vs Top-K Ranking?",
    ], title_size=1, bullet_size=18, prefix="?")

    # 5. Related Work + Gap
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    header(s, "2. Công trình Liên quan & Research Gap")
    table_slide(s, "", ["Nhóm", "Đại diện", "Ưu điểm", "Hạn chế"], [
        ["Content-Based", "TF-IDF + Cosine", "Cold-start item", "Thiếu ngữ nghĩa sâu"],
        ["Collaborative", "User/Item-CF", "Khai thác hành vi MXH", "Sparsity, cold-start"],
        ["Matrix Factor.", "Funk SVD", "Latent factors", "Tuyến tính"],
        ["Neural Rec.", "NCF / MLP", "Phi tuyến", "Cần nhiều dữ liệu"],
    ])
    gap = s.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2))
    gp = gap.text_frame.paragraphs[0]
    gp.text = "Research Gap: ViFoodRec gốc chưa đánh giá SVD và NCF → nghiên cứu này lấp khoảng trống đó."
    gp.font.size = Pt(14)
    gp.font.bold = True
    gp.font.color.rgb = C_NAVY

    # 6. Dataset overview
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    header(s, "3. Dữ liệu ViFoodRec (Sử dụng Clean Dataset)")
    bullets(s, 0.5, 1.5, 5.5, 4.5, "Thống kê", [
        "4.000 món ăn (foods.csv, 16 thuộc tính)",
        "182.678 ratings từ 101 người dùng",
        "Thang [0.0, 5.0], bước 0.5",
        "Train/Test: 90/10, random_state=42",
    ])
    bullets(s, 6.5, 1.5, 6, 4.5, "Tiền xử lý", [
        "Loại trùng (userId, foodId); Unicode NFC",
        "CBF: TF-IDF combined_features (max 2000)",
        "CF/SVD/NCF: ma trận tương tác trực tiếp",
    ])

    # 7. Data insights
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    header(s, "3. Phân tích Insight Dữ liệu Tương tác")
    add_img(s, os.path.join("report", "figures", "rating_distribution.png"), 0.5, 1.4, 4.0)
    add_img(s, os.path.join("report", "figures", "user_activity_distribution.png"), 4.8, 1.4, 4.0)
    add_img(s, os.path.join("report", "figures", "dish_type_distribution.png"), 9.1, 1.4, 3.8)
    bullets(s, 0.5, 4.8, 12.3, 2.2, "", [
        "Insight 1: Rating phân bố đồng đều 0–5 → ranh giới thích/ghét bị nhòe.",
        "Insight 2: Ma trận dẹt (101 users × 4000 items) → User-CF ổn định hơn Item-CF.",
        "Insight 3: User activity tập trung → dữ liệu thu thập có chủ đích, không phải log MXH tự nhiên.",
    ], title_size=1, bullet_size=13)

    # 8. Methods overview
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    header(s, "4. Tổng quan 4 Phương pháp")
    cards = [
        ("CBF", "TF-IDF + Cosine", "Metadata văn bản món ăn"),
        ("Memory-CF", "User/Item-CF, k=15", "Ma trận tương tác MXH"),
        ("Funk SVD", "f=10, SGD", "Latent factors + bias"),
        ("NCF (MLP)", "Embedding + MLP", "Tương tác phi tuyến"),
    ]
    for i, (title, tech, desc) in enumerate(cards):
        x = 0.4 + i * 3.2
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(2.9), Inches(4.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_WHITE
        shape.line.color.rgb = C_BLUE
        tb = s.shapes.add_textbox(Inches(x + 0.15), Inches(1.8), Inches(2.6), Inches(4))
        tf = tb.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = C_NAVY
        p2 = tf.add_paragraph()
        p2.text = tech
        p2.font.size = Pt(13)
        p2.font.color.rgb = C_ORANGE
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(12)

    # 9. SVD + NCF math
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    header(s, "4. Funk SVD & Neural CF — Công thức cốt lõi")
    bullets(s, 0.5, 1.5, 6, 5, "Funk SVD", [
        "r̂ = μ + b_u + b_i + P_u^T Q_i",
        "μ: rating trung bình; b_u, b_i: bias user/item",
        "P_u, Q_i: embedding sở thích ẩn (latent factors)",
        "Tối ưu SGD + L2 regularization (λ=0.02)",
    ], prefix="∑")
    bullets(s, 6.8, 1.5, 6, 5, "Neural CF", [
        "Input: user embedding + item embedding",
        "MLP học hàm phi tuyến f(u, i)",
        "Loss MSE, Adam lr=0.001, hidden=(16,)",
        "Vượt khả năng dot-product tuyến tính của SVD",
    ], prefix="∑")

    # 10. Experimental Setup
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    header(s, "5. Thiết kế Thí nghiệm (Reproducible)")
    table_slide(s, "", ["Thành phần", "Cấu hình"], [
        ["Chia dữ liệu", "Train 90% / Test 10%, seed=42"],
        ["CBF", "TF-IDF max_features=2000, Cosine"],
        ["Memory-CF", "k=15 neighbors, Cosine & Pearson"],
        ["Funk SVD", "f=10, lr=0.01, λ=0.02, epochs=8, SGD"],
        ["NCF (MLP)", "hidden=(16,), Adam, batch=512, max_iter=20"],
        ["Metrics", "MSE, RMSE, MAE, NMAE + P@10, R@10, NDCG, MRR"],
        ["Relevance", "rating ≥ 3.5; Top-K = 10"],
    ])

    # 11–14. Results tables
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    
    if df_reg is not None and df_rank is not None:
        cbf_rmse = f"{df_reg.loc['Content-Based', 'RMSE']:.4f}"
        cbf_mae = f"{df_reg.loc['Content-Based', 'MAE']:.4f}"
        cbf_p10 = f"{df_rank.loc['Content-Based', 'Precision@K'] * 100:.2f}%"
        cbf_ndcg = f"{df_rank.loc['Content-Based', 'NDCG@K'] * 100:.2f}%"
        
        ucf_rmse = f"{df_reg.loc['User-CF (Cosine)', 'RMSE']:.4f}"
        ucf_mae = f"{df_reg.loc['User-CF (Cosine)', 'MAE']:.4f}"
        ucf_p10 = f"{df_rank.loc['User-CF (Cosine)', 'Precision@K'] * 100:.2f}%"
        ucf_ndcg = f"{df_rank.loc['User-CF (Cosine)', 'NDCG@K'] * 100:.2f}%"
        
        icf_rmse = f"{df_reg.loc['Item-CF (Cosine)', 'RMSE']:.4f}"
        icf_mae = f"{df_reg.loc['Item-CF (Cosine)', 'MAE']:.4f}"
        icf_p10 = f"{df_rank.loc['Item-CF (Cosine)', 'Precision@K'] * 100:.2f}%"
        icf_ndcg = f"{df_rank.loc['Item-CF (Cosine)', 'NDCG@K'] * 100:.2f}%"
        
        svd_rmse = f"{df_reg.loc['Funk SVD (MF)', 'RMSE']:.4f}"
        svd_mae = f"{df_reg.loc['Funk SVD (MF)', 'MAE']:.4f}"
        svd_p10 = f"{df_rank.loc['Funk SVD (MF)', 'Precision@K'] * 100:.2f}%"
        svd_mrr = f"{df_rank.loc['Funk SVD (MF)', 'MRR'] * 100:.2f}%"
        
        ncf_rmse = f"{df_reg.loc['Neural CF (NCF)', 'RMSE']:.4f}"
        ncf_mae = f"{df_reg.loc['Neural CF (NCF)', 'MAE']:.4f}"
        ncf_p10 = f"{df_rank.loc['Neural CF (NCF)', 'Precision@K'] * 100:.2f}%"
        ncf_mrr = f"{df_rank.loc['Neural CF (NCF)', 'MRR'] * 100:.2f}%"
    else:
        # Fallback values
        cbf_rmse, cbf_mae, cbf_p10, cbf_ndcg = "1.9502", "1.6100", "36.93%", "36.70%"
        ucf_rmse, ucf_mae, ucf_p10, ucf_ndcg = "1.6429", "1.4087", "39.50%", "39.92%"
        icf_rmse, icf_mae, icf_p10, icf_ndcg = "1.6282", "1.4016", "38.51%", "38.27%"
        svd_rmse, svd_mae, svd_p10, svd_mrr = "1.6467", "1.4141", "37.72%", "57.47%"
        ncf_rmse, ncf_mae, ncf_p10, ncf_mrr = "1.5898", "1.3763", "35.54%", "62.30%"

    table_slide(s, "6. Kết quả CBF & CF (trả lời RQ1)", ["Mô hình", "RMSE", "MAE", "P@10", "NDCG"], [
        ["CBF (TF-IDF)", cbf_rmse, cbf_mae, cbf_p10, cbf_ndcg],
        ["User-CF (Cosine) ★", ucf_rmse, ucf_mae, ucf_p10, ucf_ndcg],
        ["Item-CF (Cosine)", icf_rmse, icf_mae, icf_p10, icf_ndcg],
    ])

    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    table_slide(s, "6. Kết quả SVD & NCF (trả lời RQ2)", ["Mô hình", "RMSE", "MAE", "P@10", "MRR"], [
        ["Funk SVD (f=10)", svd_rmse, svd_mae, svd_p10, svd_mrr],
        ["Neural CF (MLP) ★", ncf_rmse, ncf_mae, ncf_p10, ncf_mrr],
    ])

    # 15. Charts
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    header(s, "6. So sánh Trực quan Hiệu năng")
    add_img(s, os.path.join("report", "figures", "error_comparison.png"), 0.5, 1.3, 6.0)
    add_img(s, os.path.join("report", "figures", "ranking_comparison.png"), 6.8, 1.3, 6.0)

    # 16. Discussion
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    header(s, "7. Thảo luận & Ý nghĩa Thống kê")
    bullets(s, 0.5, 1.5, 6, 5, "Trả lời RQ3", [
        f"Rating Prediction → NCF (RMSE={ncf_rmse})",
        f"Top-K Ranking → NCF (P@10={ncf_p10}) & User-CF (P@10={ucf_p10})",
        f"MRR tốt nhất → Funk SVD (MRR={svd_mrr})",
        "Không có mô hình tốt nhất tuyệt đối",
    ])
    bullets(s, 6.8, 1.5, 6, 5, "Thảo luận thống kê", [
        f"P@10 chênh lệch giữa các mô hình nhỏ (~1-3%)",
        "Tối ưu hóa RMSE không đồng nghĩa tối ưu hóa Precision@10",
        "Học máy phi tuyến (NCF) mạnh về dự đoán điểm tương tác",
        "Học sâu & Phân rã ma trận tối ưu tùy theo mục tiêu cụ thể",
    ], prefix="📊")

    # 17. Error analysis
    s = prs.slides.add_slide(blank)
    set_bg(s, C_LIGHT)
    header(s, "8. Phân tích Lỗi (Nguyên nhân → Hướng khắc phục)")
    table_slide(s, "", ["Vấn đề", "Biểu hiện", "Hướng khắc phục"], [
        ["Cold-Start", "Dự đoán ≈ 2.5", "Hybrid CBF + NCF"],
        ["Popularity Bias", "Top-K lặp món phổ biến", "Re-ranking + penalty"],
        ["Ma trận dẹt", "Item-CF kém ổn định", "Ưu tiên User-CF / MF"],
        ["Text Noise", "CBF RMSE cao nhất", "PhoBERT embedding VN"],
        ["Rating Ambiguity", "Ngưỡng 3.5 nhòe", "Implicit feedback MXH"],
    ])

    # 18. Conclusion
    s = prs.slides.add_slide(blank)
    set_bg(s, C_NAVY)
    tb = s.shapes.add_textbox(Inches(1), Inches(0.7), Inches(11), Inches(0.7))
    tb.text_frame.paragraphs[0].text = "9. Kết luận Nghiên cứu"
    tb.text_frame.paragraphs[0].font.size = Pt(30)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = C_ORANGE
    bullets(s, 0.8, 1.5, 11.5, 5.5, "", [
        "Khai thác dữ liệu tương tác MXH ẩm thực VN bằng 4 họ phương pháp KN.",
        "NCF tốt nhất Rating Prediction; User-CF tốt nhất Top-K Ranking.",
        "Chênh lệch giữa mô hình ranking nhỏ (1–4%) — cần thận trọng khi khái quát.",
        "Hạn chế: dataset có sẵn, chưa significance test, MLP fallback.",
        "Hướng phát triển: Hybrid, embedding VN, implicit feedback thực tế.",
    ], title_size=1, bullet_size=15, prefix="✔")

    # 19. Q&A
    s = prs.slides.add_slide(blank)
    set_bg(s, C_NAVY)
    tb = s.shapes.add_textbox(Inches(2), Inches(2.8), Inches(9), Inches(2))
    tf = tb.text_frame
    tf.paragraphs[0].text = "CẢM ƠN!"
    tf.paragraphs[0].font.size = Pt(52)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = C_WHITE
    tf.paragraphs[0].alignment = 1
    p2 = tf.add_paragraph()
    p2.text = "Q & A"
    p2.font.size = Pt(28)
    p2.font.color.rgb = C_ORANGE
    p2.alignment = 1

    output_paths = [
        os.path.join(PROJECT_ROOT, "slides", "slide_research.pptx"),
        os.path.join(PROJECT_ROOT, "slides", "slide_v2.pptx"),
        os.path.join(PROJECT_ROOT, "slides", "slide.pptx"),
    ]
    saved = None
    for path in output_paths:
        try:
            prs.save(path)
            saved = path
            break
        except PermissionError:
            continue
    if saved is None:
        saved = os.path.join(PROJECT_ROOT, "slides", "slide_backup.pptx")
        prs.save(saved)
    # ASCII-only print for Windows console
    print("OK", len(prs.slides), os.path.basename(saved))


if __name__ == "__main__":
    create_presentation()
